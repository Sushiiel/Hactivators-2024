import streamlit as st
from transformers import T5Tokenizer, T5ForConditionalGeneration
from PyPDF2 import PdfReader
import cohere
import io
import re

# Set up Streamlit page configuration
st.set_page_config(page_title="Summarize, Correct, and Explain Text",
                   page_icon="📝",
                   layout='centered',
                   initial_sidebar_state='collapsed')

# Load the T5 model and tokenizer
model_name = "t5-small"
tokenizer = T5Tokenizer.from_pretrained(model_name)
model = T5ForConditionalGeneration.from_pretrained(model_name)

# Set up Cohere client (replace with your own API key)
co = cohere.Client(api_key="RQVO5wZUp59oS9eAWIr65bwT8Bw4Lu7SSpiWqfKT")

def extract_text_from_file(file):
    if file.type == "application/pdf":
        pdf_reader = PdfReader(io.BytesIO(file.read()))
        text = ""
        for page in range(len(pdf_reader.pages)):
            text += pdf_reader.pages[page].extract_text() or ""
        return text.strip()
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return extract_text_from_pptx(file)
    else:
        return file.read().decode("utf-8").strip()

def extract_text_from_pptx(file):
    from pptx import Presentation
    presentation = Presentation(io.BytesIO(file.read()))
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            # Only extract text if the shape has text and is not an image
            if hasattr(shape, "text") and shape.has_text_frame:
                text += shape.text + " "
    return text.strip()

def chunk_text(text, chunk_size=512):
    words = text.split()
    for i in range(0, len(words), chunk_size):
        yield ' '.join(words[i:i + chunk_size])

def generate_summary(text, max_words):
    input_text = "summarize: " + text
    inputs = tokenizer.encode(input_text, return_tensors="pt", max_length=1024, truncation=True)
    
    summary_ids = model.generate(
        inputs, 
        max_length=max_words,
        min_length=max_words // 2,
        length_penalty=1.0,
        num_beams=4, 
        early_stopping=True
    )

    summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
    return summary

def correct_text_with_cohere(text):
    # Pass the summarized text to the Cohere model for correction
    response = co.generate(
        model='command-xlarge',
        prompt=f"Correct the following text:\n{text}",
        max_tokens=300,
        temperature=0.5
    )
    
    # Handle the case of an empty response
    if not response.generations:
        return "No correction generated."
    
    return response.generations[0].text.strip()

def extract_and_explain_equations(text):
    # Regular expression to find equations (basic example)
    equations = re.findall(r'(\$.*?\$|\\\[(.*?)\\\])', text)
    explanations = {}
    
    for eq in equations:
        equation = eq[0]
        explanation = f"The equation {equation} represents a mathematical expression."
        # Customize your explanation logic here based on specific equations
        explanations[equation] = explanation
    
    return explanations

# Streamlit UI components
st.header("Generate, Correct Summary, and Explain Equations")

file = st.file_uploader("Upload a text, PDF, or PPTX file", type=["txt", "pdf", "pptx"])

if file is not None:
    text = extract_text_from_file(file)
    total_words = len(text.split())
    
    # Calculate half of the total word count
    half_word_count = total_words // 2
    
    st.write(f"Total number of words in the material: {total_words}")
    
    if text:
        full_summary = ""
        for chunk in chunk_text(text, 512):
            summary = generate_summary(chunk, half_word_count)  # Use half_word_count directly
            full_summary += summary + " "
        
        # Correct the summarized text with Cohere
        corrected_summary = correct_text_with_cohere(full_summary.strip())
        st.write(f"Corrected Summary (approx. {half_word_count} words): {corrected_summary}")

        # Extract and explain equations
        equations_explanations = extract_and_explain_equations(text)
        if equations_explanations:
            st.write("Equations and Explanations:")
            for eq, exp in equations_explanations.items():
                st.write(f"Equation: {eq}")
                st.write(f"Explanation: {exp}")
        else:
            st.write("No equations found in the uploaded material.")
    else:
        st.error("No text found in the uploaded file.")
else:
    st.error("Please upload a text, PDF, or PPTX file.")
